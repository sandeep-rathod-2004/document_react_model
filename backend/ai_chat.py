import os
import mimetypes
import asyncio
import base64
from fastapi import APIRouter, UploadFile, Form
from io import BytesIO
from tempfile import NamedTemporaryFile

from PIL import Image
import docx2txt
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv

# ---------------- Load Environment ----------------
load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), ".env"))
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise RuntimeError("❌ OPENAI_API_KEY not set in .env")

client = OpenAI(api_key=api_key)
router = APIRouter()

# ---------------- Helper Functions ----------------
def chunk_text(text: str, max_chars: int = 3000):
    """Split large text into chunks."""
    return [text[i:i + max_chars] for i in range(0, len(text), max_chars)]

def extract_text_from_file(filename: str, content: bytes) -> str:
    """Extract text from supported document types."""
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".txt":
            return content.decode("utf-8", errors="ignore")

        elif ext == ".pdf":
            with BytesIO(content) as f:
                reader = PdfReader(f)
                return "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext == ".docx":
            with NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                tmp.write(content)
                tmp.flush()
                text = docx2txt.process(tmp.name)
            os.remove(tmp.name)
            return text

        elif ext in [".csv", ".tsv"]:
            df = pd.read_csv(BytesIO(content))
            return df.to_string(index=False)

        elif ext == ".xlsx":
            df = pd.read_excel(BytesIO(content))
            return df.to_string(index=False)

        elif ext == ".pptx":
            prs = Presentation(BytesIO(content))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)

        else:
            return ""

    except Exception as e:
        return f"⚠️ Error extracting text: {e}"

# ---------------- OpenAI Functions ----------------
async def ask_openai_text(prompt: str):
    """Send text content to GPT-4o-mini for Q&A."""
    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": prompt},
                ],
                max_tokens=1500,
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"⚠️ OpenAI text error: {e}")
            await asyncio.sleep((attempt + 1) * 3)
    return "❌ OpenAI text request failed after multiple retries."

async def ask_openai_image(file_bytes: bytes, question: str, mime_type: str):
    """Analyze an uploaded image or PDF using GPT-4o-mini."""
    try:
        base64_data = base64.b64encode(file_bytes).decode("utf-8")
        data_url = f"data:{mime_type};base64,{base64_data}"

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": question},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }
            ],
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        return f"⚠️ Error analyzing image: {e}"

# ---------------- Unified Endpoint ----------------
@router.post("/chat")
async def chat_any(file: UploadFile, question: str = Form(...)):
    """Handle image, PDF, or document uploads and query OpenAI."""
    content = await file.read()
    mime_type, _ = mimetypes.guess_type(file.filename)
    ext = os.path.splitext(file.filename)[1].lower()

    # Handle images (png, jpg, jpeg, etc.)
    if mime_type and mime_type.startswith("image"):
        return {"answer": await ask_openai_image(content, question, mime_type), "file_type": ext}

    # Handle documents (pdf, docx, txt, csv, etc.)
    text = extract_text_from_file(file.filename, content)
    if not text.strip():
        return {"answer": "⚠️ Could not extract text from this file type.", "file_type": ext}

    chunks = chunk_text(text)
    combined_prompt = "\n\n".join([f"Chunk {i+1}:\n{chunk}" for i, chunk in enumerate(chunks)])
    final_prompt = f"""
You are a helpful assistant. Use the following document content to answer accurately.

{combined_prompt}

Question: {question}
Answer concisely:
"""
    return {"answer": await ask_openai_text(final_prompt), "file_type": ext}
